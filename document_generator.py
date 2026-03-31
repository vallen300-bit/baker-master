"""
Baker AI — Document Generator
Generates Word, Excel, PDF, and PowerPoint documents from Baker Scan output.
Called by /api/scan/generate-document endpoint in dashboard.py.
"""
import json
import logging
import os
import re
import tempfile
import uuid
from datetime import datetime

logger = logging.getLogger("baker.document_generator")

# Format-specific imports
from docx import Document
from docx.shared import Pt, Inches, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import openpyxl
from openpyxl.styles import Font, Alignment, PatternFill
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt

# Storage for generated files — maps file_id to filepath
GENERATED_FILES = {}


def generate_document(content, fmt, title, metadata=None):
    """Generate a document and store persistently in PostgreSQL."""
    generators = {
        "docx": _generate_docx,
        "xlsx": _generate_xlsx,
        "pdf": _generate_pdf,
        "pptx": _generate_pptx,
    }

    if fmt not in generators:
        raise ValueError(f"Unsupported format: {fmt}")

    file_id = str(uuid.uuid4())
    safe_title = re.sub(r'[^\w\s-]', '', title).strip().replace(' ', '_')
    date_str = datetime.utcnow().strftime('%Y-%m-%d')
    filename = f"{safe_title}_{date_str}.{fmt}"

    # Generate into temp file
    tmp_dir = tempfile.gettempdir()
    filepath = os.path.join(tmp_dir, f"baker_{file_id}.{fmt}")
    generators[fmt](content, title, filepath, metadata)

    size_bytes = os.path.getsize(filepath)

    # Read binary and store in PostgreSQL
    with open(filepath, "rb") as f:
        file_data = f.read()

    _store_in_db(file_id, filename, fmt, size_bytes, file_data, title,
                 source=metadata.get("source", "scan") if metadata else "scan")

    # NOTE: Do NOT delete the temp file here. Keep it for immediate downloads
    # via the in-memory fast path. The cleanup_old_files() function handles
    # removing stale temp files after 24h.

    # Keep in-memory cache with filepath (fast path for immediate downloads)
    GENERATED_FILES[file_id] = {
        "file_id": file_id,
        "filepath": filepath,
        "filename": filename,
        "format": fmt,
        "size_bytes": size_bytes,
        "created_at": datetime.utcnow().isoformat(),
    }

    return file_id, filename, size_bytes


def _store_in_db(file_id, filename, fmt, size_bytes, file_data, title, source="scan"):
    """Persist generated document binary to PostgreSQL."""
    try:
        from memory.store_back import SentinelStoreBack
        store = SentinelStoreBack._get_global_instance()
        conn = store._get_conn()
        if not conn:
            logger.error("No DB connection — document will only be in memory")
            return
        try:
            cur = conn.cursor()
            cur.execute("""
                INSERT INTO generated_documents (id, filename, format, size_bytes, file_data, title, source)
                VALUES (%s, %s, %s, %s, %s, %s, %s)
                ON CONFLICT (id) DO NOTHING
            """, (file_id, filename, fmt, size_bytes, file_data, title, source))
            conn.commit()
            cur.close()
            logger.info(f"Document stored in DB: {filename} ({size_bytes} bytes)")
        except Exception as e:
            conn.rollback()
            logger.error(f"Failed to store document in DB: {e}")
        finally:
            store._put_conn(conn)
    except Exception as e:
        logger.error(f"DB storage failed (non-fatal): {e}")


def get_file(file_id):
    """Retrieve file info. Check memory first, then DB."""
    # Fast path: in-memory (same server session)
    info = GENERATED_FILES.get(file_id)
    if info:
        return info

    # Slow path: load from PostgreSQL
    return _load_from_db(file_id)


def _load_from_db(file_id):
    """Load a generated document from PostgreSQL."""
    try:
        from memory.store_back import SentinelStoreBack
        store = SentinelStoreBack._get_global_instance()
        conn = store._get_conn()
        if not conn:
            return None
        try:
            cur = conn.cursor()
            cur.execute("""
                SELECT filename, format, size_bytes, file_data
                FROM generated_documents
                WHERE id = %s AND expired = FALSE
            """, (file_id,))
            row = cur.fetchone()
            cur.close()
            if not row:
                return None

            # Write binary back to temp file for FileResponse
            tmp_dir = tempfile.gettempdir()
            filepath = os.path.join(tmp_dir, f"baker_{file_id}.{row[1]}")
            with open(filepath, "wb") as f:
                f.write(row[3])  # file_data (bytes)

            info = {
                "file_id": file_id,
                "filepath": filepath,
                "filename": row[0],
                "format": row[1],
                "size_bytes": row[2],
            }

            # Cache in memory for future requests
            GENERATED_FILES[file_id] = info

            # Update downloaded_at
            try:
                cur2 = conn.cursor()
                cur2.execute("""
                    UPDATE generated_documents SET downloaded_at = NOW() WHERE id = %s
                """, (file_id,))
                conn.commit()
                cur2.close()
            except Exception:
                conn.rollback()

            return info
        except Exception as e:
            conn.rollback()
            logger.error(f"Failed to load document from DB: {e}")
            return None
        finally:
            store._put_conn(conn)
    except Exception:
        return None


def list_generated_documents(limit=20):
    """List recent generated documents (for the right panel)."""
    try:
        from memory.store_back import SentinelStoreBack
        store = SentinelStoreBack._get_global_instance()
        conn = store._get_conn()
        if not conn:
            return []
        try:
            cur = conn.cursor()
            cur.execute("""
                SELECT id, filename, format, size_bytes, title, source, created_at
                FROM generated_documents
                WHERE expired = FALSE
                ORDER BY created_at DESC
                LIMIT %s
            """, (limit,))
            rows = cur.fetchall()
            cur.close()
            return [
                {
                    "file_id": r[0],
                    "filename": r[1],
                    "format": r[2],
                    "size_bytes": r[3],
                    "title": r[4],
                    "source": r[5],
                    "created_at": r[6].isoformat() if r[6] else None,
                    "download_url": f"/api/scan/download/{r[0]}",
                }
                for r in rows
            ]
        except Exception as e:
            conn.rollback()
            logger.error(f"Failed to list generated documents: {e}")
            return []
        finally:
            store._put_conn(conn)
    except Exception:
        return []


def cleanup_old_files(max_age_days=30):
    """Expire documents older than max_age_days in PostgreSQL. Call from scheduler."""
    try:
        from memory.store_back import SentinelStoreBack
        store = SentinelStoreBack._get_global_instance()
        conn = store._get_conn()
        if not conn:
            return
        try:
            cur = conn.cursor()
            cur.execute("""
                UPDATE generated_documents
                SET expired = TRUE
                WHERE created_at < NOW() - make_interval(days => %s)
                  AND expired = FALSE
            """, (max_age_days,))
            count = cur.rowcount
            conn.commit()
            cur.close()
            if count:
                logger.info(f"Expired {count} old generated documents")
        except Exception as e:
            conn.rollback()
            logger.error(f"Document cleanup failed: {e}")
        finally:
            store._put_conn(conn)
    except Exception:
        pass

    # Also clean in-memory cache (still useful for temp files)
    now = datetime.utcnow()
    expired = []
    for fid, info in GENERATED_FILES.items():
        created_str = info.get("created_at")
        if created_str:
            try:
                created = datetime.fromisoformat(created_str)
                if (now - created).total_seconds() > 86400:  # 24h for in-memory
                    expired.append(fid)
            except Exception:
                pass
    for fid in expired:
        info = GENERATED_FILES.pop(fid, None)
        if info and info.get("filepath"):
            try:
                os.remove(info["filepath"])
            except OSError:
                pass


# ============================================================
# Generator: Word (.docx)
# ============================================================

def _generate_docx(content, title, filepath, metadata):
    """Parse markdown content -> Word document with clean business formatting."""
    doc = Document()

    # Title
    heading = doc.add_heading(title, level=0)
    heading.alignment = WD_ALIGN_PARAGRAPH.LEFT

    # Metadata line
    if metadata:
        meta_para = doc.add_paragraph()
        meta_run = meta_para.add_run(
            f"Generated by {metadata.get('generated_by', 'Baker Scan')} — "
            f"{metadata.get('timestamp', datetime.utcnow().isoformat())}"
        )
        meta_run.font.size = Pt(9)
        meta_run.font.color.rgb = RGBColor(128, 128, 128)
        doc.add_paragraph()  # spacer

    # Parse markdown lines into document elements
    lines = content.split('\n')
    for line in lines:
        stripped = line.strip()
        if not stripped:
            doc.add_paragraph()
            continue

        # Headings
        if stripped.startswith('### '):
            doc.add_heading(stripped[4:], level=3)
        elif stripped.startswith('## '):
            doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith('# '):
            doc.add_heading(stripped[2:], level=1)
        # Bullet points
        elif stripped.startswith('- ') or stripped.startswith('* '):
            doc.add_paragraph(stripped[2:], style='List Bullet')
        # Numbered items
        elif re.match(r'^\d+\.\s', stripped):
            text = re.sub(r'^\d+\.\s', '', stripped)
            doc.add_paragraph(text, style='List Number')
        # Regular paragraph
        else:
            para = doc.add_paragraph()
            # Handle basic bold (**text**) and italic (*text*)
            parts = re.split(r'(\*\*.*?\*\*|\*.*?\*)', stripped)
            for part in parts:
                if part.startswith('**') and part.endswith('**'):
                    run = para.add_run(part[2:-2])
                    run.bold = True
                elif part.startswith('*') and part.endswith('*'):
                    run = para.add_run(part[1:-1])
                    run.italic = True
                else:
                    para.add_run(part)

    doc.save(filepath)


# ============================================================
# Generator: Excel (.xlsx)
# ============================================================

def _generate_xlsx(content, title, filepath, metadata):
    """
    Content should be JSON-parseable with structure:
    { "headers": ["Col A", "Col B"], "rows": [["val1", "val2"], ...] }
    Falls back to parsing markdown tables if JSON parsing fails.
    """
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = title[:31]  # Excel sheet name max 31 chars

    header_font = Font(bold=True, size=11, color="FFFFFF")
    header_fill = PatternFill(start_color="2C3E50", end_color="2C3E50", fill_type="solid")

    try:
        data = json.loads(content) if isinstance(content, str) else content
        headers = data.get("headers", [])
        rows = data.get("rows", [])
    except (json.JSONDecodeError, AttributeError):
        # Fallback: parse markdown table
        headers, rows = _parse_markdown_table(content)

    # Write headers
    for col_idx, header in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col_idx, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal='center')

    # Write data rows
    for row_idx, row_data in enumerate(rows, 2):
        for col_idx, value in enumerate(row_data, 1):
            cell = ws.cell(row=row_idx, column=col_idx, value=value)
            cell.alignment = Alignment(wrap_text=True)

    # Auto-width columns
    for col in ws.columns:
        max_length = 0
        col_letter = col[0].column_letter
        for cell in col:
            if cell.value:
                max_length = max(max_length, len(str(cell.value)))
        ws.column_dimensions[col_letter].width = min(max_length + 4, 50)

    # Metadata row at bottom
    if metadata:
        skip_row = len(rows) + 3
        ws.cell(row=skip_row, column=1,
                value=f"Generated by {metadata.get('generated_by', 'Baker Scan')} — {metadata.get('timestamp', '')}")
        ws.cell(row=skip_row, column=1).font = Font(size=8, color="999999")

    wb.save(filepath)


def _parse_markdown_table(text):
    """Extract headers and rows from a markdown table string."""
    lines = [l.strip() for l in text.strip().split('\n') if l.strip()]
    headers = []
    rows = []
    for line in lines:
        if '|' not in line:
            continue
        cells = [c.strip() for c in line.split('|') if c.strip()]
        # Skip separator rows (---|---)
        if all(re.match(r'^[-:]+$', c) for c in cells):
            continue
        if not headers:
            headers = cells
        else:
            rows.append(cells)
    return headers, rows


# ============================================================
# Generator: PDF (.pdf)
# ============================================================

def _generate_pdf(content, title, filepath, metadata):
    """Markdown content -> clean PDF via reportlab."""
    doc = SimpleDocTemplate(
        filepath, pagesize=A4,
        leftMargin=inch, rightMargin=inch,
        topMargin=inch, bottomMargin=inch
    )

    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle(
        'BakerTitle', parent=styles['Title'],
        fontSize=18, spaceAfter=12,
        textColor=colors.HexColor('#1a1a2e')
    )
    heading_style = ParagraphStyle(
        'BakerH2', parent=styles['Heading2'],
        fontSize=14, spaceBefore=16, spaceAfter=8,
        textColor=colors.HexColor('#2C3E50')
    )
    body_style = ParagraphStyle(
        'BakerBody', parent=styles['Normal'],
        fontSize=10, leading=14, spaceAfter=6
    )
    meta_style = ParagraphStyle(
        'BakerMeta', parent=styles['Normal'],
        fontSize=8, textColor=colors.HexColor('#999999')
    )

    elements = []

    # Title
    elements.append(Paragraph(title, title_style))

    # Metadata
    if metadata:
        meta_text = (f"Generated by {metadata.get('generated_by', 'Baker Scan')} — "
                     f"{metadata.get('timestamp', '')}")
        elements.append(Paragraph(meta_text, meta_style))

    elements.append(Spacer(1, 12))

    # Parse markdown -> reportlab elements
    lines = content.split('\n')
    for line in lines:
        stripped = line.strip()
        if not stripped:
            elements.append(Spacer(1, 6))
            continue

        if stripped.startswith('## '):
            elements.append(Paragraph(stripped[3:], heading_style))
        elif stripped.startswith('# '):
            elements.append(Paragraph(stripped[2:], title_style))
        elif stripped.startswith('- ') or stripped.startswith('* '):
            bullet_text = f"\u2022 {stripped[2:]}"
            elements.append(Paragraph(bullet_text, body_style))
        else:
            # Convert **bold** to <b>bold</b> for reportlab
            text = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', stripped)
            text = re.sub(r'\*(.*?)\*', r'<i>\1</i>', text)
            elements.append(Paragraph(text, body_style))

    doc.build(elements)


# ============================================================
# Generator: PowerPoint (.pptx)
# ============================================================

def _generate_pptx(content, title, filepath, metadata):
    """
    Content should be JSON-parseable with structure:
    { "slides": [{"title": "...", "bullets": ["...", ...]}, ...] }
    Falls back to splitting by ## headings if JSON parsing fails.
    """
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    try:
        data = json.loads(content) if isinstance(content, str) else content
        slides = data.get("slides", [])
    except (json.JSONDecodeError, AttributeError):
        slides = _parse_markdown_slides(content)

    # Title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        meta_text = "Generated by Baker Scan"
        if metadata and metadata.get('timestamp'):
            meta_text += f" — {metadata['timestamp']}"
        slide.placeholders[1].text = meta_text

    # Content slides
    content_layout = prs.slide_layouts[1]  # Title + Content
    for slide_data in slides:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = slide_data.get("title", "")

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        bullets = slide_data.get("bullets", [])
        for i, bullet in enumerate(bullets):
            if i == 0:
                tf.paragraphs[0].text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet
            tf.paragraphs[-1].font.size = PptxPt(18)

    prs.save(filepath)


def _parse_markdown_slides(text):
    """Split markdown by ## headings into slide structures."""
    slides = []
    current_slide = None

    for line in text.split('\n'):
        stripped = line.strip()
        if stripped.startswith('## '):
            if current_slide:
                slides.append(current_slide)
            current_slide = {"title": stripped[3:], "bullets": []}
        elif current_slide and (stripped.startswith('- ') or stripped.startswith('* ')):
            current_slide["bullets"].append(stripped[2:])
        elif current_slide and stripped and not stripped.startswith('#'):
            current_slide["bullets"].append(stripped)

    if current_slide:
        slides.append(current_slide)

    return slides


# ============================================================
# Professional Dossier DOCX Generator (ART-1)
# ============================================================

# Colors
_NAVY = RGBColor(0x1a, 0x1a, 0x2e)
_DARK_BLUE = RGBColor(0x2C, 0x3E, 0x50)
_BLUE_ACCENT = RGBColor(0x0a, 0x6f, 0xdb)
_GRAY = RGBColor(0x80, 0x80, 0x80)
_LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _add_thin_border(paragraph, color="1a1a2e"):
    """Add a thin bottom border to a paragraph (used as horizontal rule)."""
    pPr = paragraph._p.get_or_add_pPr()
    pBdr = OxmlElement('w:pBdr')
    bottom = OxmlElement('w:bottom')
    bottom.set(qn('w:val'), 'single')
    bottom.set(qn('w:sz'), '4')
    bottom.set(qn('w:space'), '4')
    bottom.set(qn('w:color'), color)
    pBdr.append(bottom)
    pPr.append(pBdr)


def _set_cell_shading(cell, color_hex):
    """Set background shading on a table cell."""
    shading = OxmlElement('w:shd')
    shading.set(qn('w:fill'), color_hex)
    shading.set(qn('w:val'), 'clear')
    cell._tc.get_or_add_tcPr().append(shading)


def _add_header_footer(doc, title_text="Baker Research Dossier"):
    """Add header and footer to the document."""
    for section in doc.sections:
        # Header
        header = section.header
        header.is_linked_to_previous = False
        hp = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
        hp.text = ""
        run = hp.add_run(title_text)
        run.font.size = Pt(8)
        run.font.color.rgb = _GRAY
        hp.alignment = WD_ALIGN_PARAGRAPH.LEFT
        _add_thin_border(hp, "CCCCCC")

        # Footer with page number
        footer = section.footer
        footer.is_linked_to_previous = False
        fp = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
        fp.text = ""
        fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = fp.add_run("Page ")
        run.font.size = Pt(8)
        run.font.color.rgb = _GRAY
        # Page number field
        fld_char1 = OxmlElement('w:fldChar')
        fld_char1.set(qn('w:fldCharType'), 'begin')
        run._r.append(fld_char1)
        instr = OxmlElement('w:instrText')
        instr.set(qn('xml:space'), 'preserve')
        instr.text = ' PAGE '
        run._r.append(instr)
        fld_char2 = OxmlElement('w:fldChar')
        fld_char2.set(qn('w:fldCharType'), 'end')
        run._r.append(fld_char2)


def _add_cover_page(doc, subject_name, subject_type, specialists_text, date_str):
    """Add a professional cover page."""
    # Top spacer
    for _ in range(6):
        doc.add_paragraph()

    # Title
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title_para.add_run("RESEARCH DOSSIER")
    run.font.size = Pt(28)
    run.font.color.rgb = _NAVY
    run.bold = True
    run.font.name = "Calibri"

    doc.add_paragraph()

    # Subject name
    subj_para = doc.add_paragraph()
    subj_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subj_para.add_run(subject_name)
    run.font.size = Pt(22)
    run.font.color.rgb = _DARK_BLUE
    run.font.name = "Calibri"

    # Subject type
    type_para = doc.add_paragraph()
    type_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = type_para.add_run(subject_type.upper())
    run.font.size = Pt(11)
    run.font.color.rgb = _GRAY
    run.font.name = "Calibri"

    # Spacer
    for _ in range(4):
        doc.add_paragraph()

    # Thin rule
    rule = doc.add_paragraph()
    rule.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _add_thin_border(rule, "CCCCCC")

    # Metadata block
    meta_table = doc.add_table(rows=4, cols=2)
    meta_table.autofit = True
    meta_table.columns[0].width = Inches(2)
    meta_table.columns[1].width = Inches(4)

    labels = ["Date", "Classification", "Specialists", "Prepared by"]
    values = [date_str, "CONFIDENTIAL", specialists_text,
              "Baker Research Engine"]

    for i, (label, value) in enumerate(zip(labels, values)):
        left = meta_table.cell(i, 0)
        right = meta_table.cell(i, 1)
        lp = left.paragraphs[0]
        rp = right.paragraphs[0]
        lr = lp.add_run(label)
        lr.font.size = Pt(10)
        lr.font.color.rgb = _GRAY
        lr.font.name = "Calibri"
        rr = rp.add_run(value)
        rr.font.size = Pt(10)
        rr.font.color.rgb = _NAVY
        rr.font.name = "Calibri"
        if label == "Classification":
            rr.bold = True

    # Remove table borders
    for row in meta_table.rows:
        for cell in row.cells:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            tcBorders = OxmlElement('w:tcBorders')
            for edge in ('top', 'left', 'bottom', 'right'):
                element = OxmlElement(f'w:{edge}')
                element.set(qn('w:val'), 'none')
                element.set(qn('w:sz'), '0')
                element.set(qn('w:space'), '0')
                element.set(qn('w:color'), 'auto')
                tcBorders.append(element)
            tcPr.append(tcBorders)

    # Page break after cover
    doc.add_page_break()


def _render_markdown_section(doc, content):
    """Render markdown content into Word document elements with professional formatting."""
    lines = content.split('\n')
    for line in lines:
        stripped = line.strip()

        if not stripped:
            # Skip blank lines — paragraph spacing handles visual gaps
            continue

        # Horizontal rule → thin border paragraph
        if stripped in ('---', '***', '___'):
            rule = doc.add_paragraph()
            _add_thin_border(rule, "CCCCCC")
            continue

        # Headings
        if stripped.startswith('#### '):
            h = doc.add_heading(stripped[5:], level=4)
            for run in h.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(11)
                run.font.color.rgb = _DARK_BLUE
            continue
        if stripped.startswith('### '):
            h = doc.add_heading(stripped[4:], level=3)
            for run in h.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(12)
                run.font.color.rgb = _DARK_BLUE
            continue
        if stripped.startswith('## '):
            # This is a specialist section header — handled by caller
            h = doc.add_heading(stripped[3:], level=2)
            for run in h.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(14)
                run.font.color.rgb = _NAVY
                run.bold = True
            continue
        if stripped.startswith('# '):
            h = doc.add_heading(stripped[2:], level=1)
            for run in h.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(16)
                run.font.color.rgb = _NAVY
            continue

        # Bullet points
        if stripped.startswith('- ') or stripped.startswith('* '):
            p = doc.add_paragraph(style='List Bullet')
            _add_inline_formatting(p, stripped[2:])
            for run in p.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(11)
            continue

        # Sub-bullets
        if stripped.startswith('  - ') or stripped.startswith('  * '):
            p = doc.add_paragraph(style='List Bullet 2')
            _add_inline_formatting(p, stripped[4:])
            for run in p.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(10)
            continue

        # Numbered items
        if re.match(r'^\d+\.\s', stripped):
            text = re.sub(r'^\d+\.\s', '', stripped)
            p = doc.add_paragraph(style='List Number')
            _add_inline_formatting(p, text)
            for run in p.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(11)
            continue

        # Regular paragraph
        para = doc.add_paragraph()
        _add_inline_formatting(para, stripped)
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(11)


def _add_inline_formatting(paragraph, text):
    """Parse **bold** and *italic* inline markdown and add runs to paragraph."""
    parts = re.split(r'(\*\*.*?\*\*|\*.*?\*)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            run = paragraph.add_run(part[2:-2])
            run.bold = True
        elif part.startswith('*') and part.endswith('*'):
            run = paragraph.add_run(part[1:-1])
            run.italic = True
        else:
            paragraph.add_run(part)


def generate_dossier_docx(dossier_md, subject_name, subject_type,
                          specialists_text, filepath):
    """
    Generate a professional McKinsey-style research dossier .docx.

    This splits the markdown by ## sections (specialist headers) and adds:
    - Cover page with metadata
    - Page breaks between specialist sections
    - Professional fonts (Calibri) and consistent styling
    - Header/footer with page numbers
    - Horizontal rules as thin borders

    Args:
        dossier_md: Full dossier markdown content
        subject_name: Name of the research subject
        subject_type: Type (person, company, etc.)
        specialists_text: Comma-separated specialist names
        filepath: Output .docx file path
    """
    doc = Document()

    # Set default font + paragraph spacing for readability
    style = doc.styles['Normal']
    font = style.font
    font.name = "Calibri"
    font.size = Pt(11)
    font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    # Line spacing 1.15 + 6pt after paragraph — much more readable
    style.paragraph_format.line_spacing = 1.15
    style.paragraph_format.space_after = Pt(6)
    style.paragraph_format.space_before = Pt(0)

    # Heading styles — add spacing above for visual separation
    for lvl in (1, 2, 3, 4):
        try:
            hstyle = doc.styles[f'Heading {lvl}']
            hstyle.paragraph_format.space_before = Pt(18 if lvl <= 2 else 12)
            hstyle.paragraph_format.space_after = Pt(6)
        except Exception:
            pass

    # List styles — tighter spacing
    for lst in ('List Bullet', 'List Bullet 2', 'List Number'):
        try:
            lstyle = doc.styles[lst]
            lstyle.paragraph_format.space_after = Pt(3)
            lstyle.paragraph_format.space_before = Pt(1)
            lstyle.paragraph_format.line_spacing = 1.15
        except Exception:
            pass

    # Set narrow margins
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.0)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    date_str = datetime.utcnow().strftime("%d %B %Y")

    # 1. Cover page
    _add_cover_page(doc, subject_name, subject_type, specialists_text, date_str)

    # 2. Header/footer (applied to all sections)
    _add_header_footer(doc)

    # 3. Split content by ## sections and render with page breaks
    # Skip the cover-page markdown (# title, metadata lines, first ---)
    sections = re.split(r'\n(?=## )', dossier_md)
    is_first_section = True

    for section_md in sections:
        section_md = section_md.strip()
        if not section_md:
            continue

        # Skip the top-level header block (# Research Dossier: ...)
        if section_md.startswith('# ') and not section_md.startswith('## '):
            # This is the metadata header — skip (cover page handles it)
            continue

        # Page break between specialist sections (not before the first)
        if section_md.startswith('## ') and not is_first_section:
            doc.add_page_break()

        if section_md.startswith('## '):
            is_first_section = False

        # Skip the footer line
        if section_md.startswith('*Generated by Baker'):
            continue

        _render_markdown_section(doc, section_md)

    # Final footer note
    doc.add_paragraph()
    rule = doc.add_paragraph()
    _add_thin_border(rule, "CCCCCC")
    footer_para = doc.add_paragraph()
    footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = footer_para.add_run(
        f"Generated by Baker Research Engine  |  {date_str}  |  CONFIDENTIAL"
    )
    run.font.size = Pt(8)
    run.font.color.rgb = _GRAY
    run.font.name = "Calibri"

    doc.save(filepath)


# ============================================================
# Brisen External Document Generator
# ============================================================
# For external-facing documents: legal memos, PR strategies,
# investor materials, dossiers. Matches 007 Brisen house style.
#
# Typography: Arial throughout (headings + body)
# Palette: Brisen dark blue #1B3A5C, body #333333, meta #666666
# Header: "BRISEN GROUP" + "CONFIDENTIAL — {type}"
# Footer: "{doc#} — {title}" + "Page X"
# End block: "Prepared {date} | Brisen Group Strategic Intelligence Unit"
# NO mention of Baker, AI, or automation anywhere.
# ============================================================

# Brisen palette (matched from 007 document)
_BRS_BLUE = RGBColor(0x1B, 0x3A, 0x5C)        # Primary headings + header
_BRS_BODY = RGBColor(0x33, 0x33, 0x33)         # Body text
_BRS_META = RGBColor(0x66, 0x66, 0x66)         # Metadata, footer, cover subtitle
_BRS_WHITE = RGBColor(0xFF, 0xFF, 0xFF)        # Table header text
_BRS_ORANGE = RGBColor(0xE0, 0x70, 0x20)       # CRITICAL callout title
_BRS_RED = RGBColor(0xC0, 0x00, 0x00)          # PROHIBITION callout title
_BRS_TABLE_HDR = "1B3A5C"                       # Table header bg
_BRS_TABLE_ALT = "F2F2F2"                       # Alternating row bg
_BRS_EXEC_BG = "EDF2F7"                         # Exec summary box bg
_BRS_CALLOUT_BG = "F5F5F5"                      # Generic callout bg


def generate_mckinsey_docx(
    content,
    title,
    subtitle="",
    classification="CONFIDENTIAL",
    filepath=None,
    prepared_by="Brisen Group — Strategic Intelligence Unit",
    prepared_for="",
    date_str=None,
    exec_summary=None,
    doc_number="",
    doc_type="",
    supersedes="",
):
    """
    Generate a Brisen-style external document (.docx).
    Matches the 007 house style exactly: all Arial, #1B3A5C palette,
    BRISEN GROUP header, document-numbered footer, callout boxes.

    Args:
        content: str — markdown body content
        title: str — main title on cover (e.g., "HAGENAUER INSOLVENCY")
        subtitle: str — subtitle line (e.g., "PR & PRESS STRATEGY — PART 2")
        classification: str — "CONFIDENTIAL", "STRICTLY CONFIDENTIAL", etc.
        filepath: str — output path. Auto-generated if None.
        prepared_by: str — attribution line
        prepared_for: str — recipient (e.g., "Sandra Luger, Gaisberg Consulting")
        date_str: str — date string. Auto-generated if None.
        exec_summary: str — optional executive summary (rendered in shaded box)
        doc_number: str — document number for footer (e.g., "004")
        doc_type: str — header right text (e.g., "PR STRATEGY")
        supersedes: str — optional supersedes line (e.g., "v1.0 (16 March 2026)")

    Returns:
        filepath (str) — path to generated .docx
    """
    if filepath is None:
        safe = re.sub(r'[^\w\s-]', '', title).strip().replace(' ', '_')
        filepath = os.path.join(
            tempfile.gettempdir(),
            f"brisen_{safe}_{datetime.utcnow().strftime('%Y%m%d')}.docx"
        )

    if date_str is None:
        date_str = datetime.utcnow().strftime("%d %B %Y")

    doc = Document()

    # -- Global styles --
    _brs_setup_styles(doc)

    # -- US Letter page, 2.5cm margins --
    for section in doc.sections:
        section.page_width = Inches(8.5)
        section.page_height = Inches(11)
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    # 1. Cover page
    _brs_cover_page(doc, title, subtitle, classification, prepared_by,
                    prepared_for, date_str, supersedes)

    # 2. Header + footer
    _brs_header_footer(doc, doc_number, title, doc_type, classification)

    # 3. Executive summary box (if provided)
    if exec_summary:
        _brs_exec_summary_box(doc, exec_summary)

    # 4. Body content
    _brs_render_body(doc, content)

    # 5. End-of-document block
    doc.add_paragraph()
    ep = doc.add_paragraph()
    ep.alignment = WD_ALIGN_PARAGRAPH.CENTER
    er = ep.add_run(f"END OF DOCUMENT{' ' + doc_number if doc_number else ''}")
    er.font.size = Pt(10)
    er.font.color.rgb = _BRS_META
    er.font.name = "Arial"
    er.bold = True

    pp = doc.add_paragraph()
    pp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    pr = pp.add_run(f"Prepared {date_str} | {prepared_by}")
    pr.font.size = Pt(8)
    pr.font.color.rgb = _BRS_META
    pr.font.name = "Arial"

    if supersedes:
        sp = doc.add_paragraph()
        sp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sr = sp.add_run(f"Supersedes {supersedes}")
        sr.font.size = Pt(7)
        sr.font.color.rgb = _BRS_META
        sr.font.name = "Arial"
        sr.italic = True

    doc.save(filepath)
    return filepath


def generate_internal_docx(content, title, filepath=None, date_str=None):
    """
    Generate a clean internal Word document (.docx).
    Simple, readable, no fancy formatting. For Director/team use.
    """
    if filepath is None:
        safe = re.sub(r'[^\w\s-]', '', title).strip().replace(' ', '_')
        filepath = os.path.join(
            tempfile.gettempdir(),
            f"internal_{safe}_{datetime.utcnow().strftime('%Y%m%d')}.docx"
        )

    if date_str is None:
        date_str = datetime.utcnow().strftime("%d %B %Y")

    doc = Document()

    style = doc.styles['Normal']
    style.font.name = "Calibri"
    style.font.size = Pt(11)
    style.font.color.rgb = _BRS_BODY
    style.paragraph_format.line_spacing = 1.15
    style.paragraph_format.space_after = Pt(6)

    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.0)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    h = doc.add_heading(title, level=0)
    h.alignment = WD_ALIGN_PARAGRAPH.LEFT
    for run in h.runs:
        run.font.name = "Calibri"
        run.font.color.rgb = _BRS_BLUE

    dp = doc.add_paragraph()
    dr = dp.add_run(date_str)
    dr.font.size = Pt(10)
    dr.font.color.rgb = _BRS_META
    dr.font.name = "Calibri"
    doc.add_paragraph()

    _brs_render_body(doc, content, font_name="Calibri")

    doc.save(filepath)
    return filepath


# --- Brisen house style internals ---

def _brs_setup_styles(doc):
    """Configure global styles matching 007 Brisen house style."""
    style = doc.styles['Normal']
    style.font.name = "Arial"
    style.font.size = Pt(10)
    style.font.color.rgb = _BRS_BODY
    style.paragraph_format.line_spacing = 1.15
    style.paragraph_format.space_after = Pt(6)
    style.paragraph_format.space_before = Pt(0)

    # All headings: Arial, #1B3A5C
    for lvl, (size, bold) in {
        1: (14, True),
        2: (12, True),
        3: (11, True),
        4: (10, False),
    }.items():
        try:
            hs = doc.styles[f'Heading {lvl}']
            hs.font.name = "Arial"
            hs.font.size = Pt(size)
            hs.font.color.rgb = _BRS_BLUE
            hs.font.bold = bold
            hs.paragraph_format.space_before = Pt(18 if lvl <= 2 else 12)
            hs.paragraph_format.space_after = Pt(6)
        except Exception:
            pass

    for lst in ('List Bullet', 'List Bullet 2', 'List Number', 'List Paragraph'):
        try:
            ls = doc.styles[lst]
            ls.font.name = "Arial"
            ls.font.size = Pt(10)
            ls.paragraph_format.space_after = Pt(3)
            ls.paragraph_format.space_before = Pt(1)
            ls.paragraph_format.line_spacing = 1.15
        except Exception:
            pass


def _brs_cover_page(doc, title, subtitle, classification, prepared_by,
                    prepared_for, date_str, supersedes):
    """Brisen cover page — matches 007 exactly."""
    # Top spacer
    for _ in range(3):
        doc.add_paragraph()

    # BRISEN GROUP
    bg = doc.add_paragraph()
    bg.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = bg.add_run("BRISEN GROUP")
    r.font.name = "Arial"
    r.font.size = Pt(18)
    r.font.color.rgb = _BRS_BLUE
    r.bold = True

    # Classification subtitle
    cs = doc.add_paragraph()
    cs.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = cs.add_run(f"Confidential {classification.title().replace('Confidential', '').strip() or 'Communications'}".strip())
    r.font.name = "Arial"
    r.font.size = Pt(11)
    r.font.color.rgb = _BRS_META

    doc.add_paragraph()

    # Main title
    tp = doc.add_paragraph()
    tp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = tp.add_run(title.upper())
    r.font.name = "Arial"
    r.font.size = Pt(22)
    r.font.color.rgb = _BRS_BLUE
    r.bold = True

    # Subtitle
    if subtitle:
        sp = doc.add_paragraph()
        sp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = sp.add_run(subtitle)
        r.font.name = "Arial"
        r.font.size = Pt(16)
        r.font.color.rgb = _BRS_BODY
        r.bold = True

    # Prepared for
    if prepared_for:
        doc.add_paragraph()
        pf = doc.add_paragraph()
        pf.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = pf.add_run(f"Prepared for {prepared_for}")
        r.font.name = "Arial"
        r.font.size = Pt(13)
        r.font.color.rgb = _BRS_META

    # Spacer
    for _ in range(4):
        doc.add_paragraph()

    # Metadata table — borderless
    n_rows = 3  # Date, Classification, Prepared by
    if prepared_for:
        n_rows += 1
    if supersedes:
        n_rows += 1

    meta = doc.add_table(rows=n_rows, cols=2)
    meta.autofit = True
    meta.columns[0].width = Inches(1.8)
    meta.columns[1].width = Inches(4.2)

    rows_data = [("Date:", date_str), ("Classification:", classification)]
    if prepared_for:
        rows_data.append(("Prepared for:", prepared_for))
    rows_data.append(("Prepared by:", prepared_by))
    if supersedes:
        rows_data.append(("Supersedes:", supersedes))

    for i, (label, value) in enumerate(rows_data):
        left = meta.cell(i, 0)
        right = meta.cell(i, 1)
        lp = left.paragraphs[0]
        rp = right.paragraphs[0]
        lr = lp.add_run(label)
        lr.font.name = "Arial"
        lr.font.size = Pt(10)
        lr.font.color.rgb = _BRS_META
        lr.bold = True
        rr = rp.add_run(value)
        rr.font.name = "Arial"
        rr.font.size = Pt(10)
        rr.font.color.rgb = _BRS_BODY

    # Remove table borders
    for row in meta.rows:
        for cell in row.cells:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            tcBorders = OxmlElement('w:tcBorders')
            for edge in ('top', 'left', 'bottom', 'right'):
                el = OxmlElement(f'w:{edge}')
                el.set(qn('w:val'), 'none')
                el.set(qn('w:sz'), '0')
                el.set(qn('w:space'), '0')
                el.set(qn('w:color'), 'auto')
                tcBorders.append(el)
            tcPr.append(tcBorders)

    doc.add_page_break()


def _brs_header_footer(doc, doc_number, title, doc_type, classification):
    """Header: BRISEN GROUP + CONFIDENTIAL — TYPE. Footer: doc# + title + Page."""
    for section in doc.sections:
        # Header: "BRISEN GROUP\tCONFIDENTIAL — TYPE"
        header = section.header
        header.is_linked_to_previous = False
        hp = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
        hp.text = ""

        r1 = hp.add_run("BRISEN GROUP")
        r1.font.name = "Arial"
        r1.font.size = Pt(8)
        r1.font.color.rgb = _BRS_BLUE
        r1.bold = True

        hp.add_run("\t")  # tab to right

        right_text = f"{classification} — {doc_type}" if doc_type else classification
        r2 = hp.add_run(right_text)
        r2.font.name = "Arial"
        r2.font.size = Pt(7)
        r2.font.color.rgb = _BRS_META

        _add_thin_border(hp, "CCCCCC")

        # Footer: "doc# — title\tPage X"
        footer = section.footer
        footer.is_linked_to_previous = False
        fp = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
        fp.text = ""

        # Build short title for footer
        footer_title = title[:50]
        if doc_number:
            footer_left = f"{doc_number} — {footer_title}"
        else:
            footer_left = footer_title

        r1 = fp.add_run(footer_left)
        r1.font.name = "Arial"
        r1.font.size = Pt(7)
        r1.font.color.rgb = _BRS_META

        fp.add_run("\t")

        r2 = fp.add_run("Page ")
        r2.font.name = "Arial"
        r2.font.size = Pt(7)
        r2.font.color.rgb = _BRS_META
        fld1 = OxmlElement('w:fldChar')
        fld1.set(qn('w:fldCharType'), 'begin')
        r2._r.append(fld1)
        instr = OxmlElement('w:instrText')
        instr.set(qn('xml:space'), 'preserve')
        instr.text = ' PAGE '
        r2._r.append(instr)
        fld2 = OxmlElement('w:fldChar')
        fld2.set(qn('w:fldCharType'), 'end')
        r2._r.append(fld2)


def _brs_exec_summary_box(doc, summary_text):
    """Executive summary in shaded box with Brisen blue border."""
    h = doc.add_heading("Executive Summary", level=1)
    for run in h.runs:
        run.font.name = "Arial"
        run.font.size = Pt(14)
        run.font.color.rgb = _BRS_BLUE

    table = doc.add_table(rows=1, cols=1)
    table.autofit = True
    cell = table.cell(0, 0)
    _set_cell_shading(cell, _BRS_EXEC_BG)

    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcMar = OxmlElement('w:tcMar')
    for edge in ('top', 'bottom', 'start', 'end'):
        mar = OxmlElement(f'w:{edge}')
        mar.set(qn('w:w'), '180')
        mar.set(qn('w:type'), 'dxa')
        tcMar.append(mar)
    tcPr.append(tcMar)

    tcBorders = OxmlElement('w:tcBorders')
    for edge in ('top', 'left', 'bottom', 'right'):
        el = OxmlElement(f'w:{edge}')
        el.set(qn('w:val'), 'single')
        el.set(qn('w:sz'), '4')
        el.set(qn('w:space'), '0')
        el.set(qn('w:color'), '1B3A5C')
        tcBorders.append(el)
    tcPr.append(tcBorders)

    lines = summary_text.strip().split('\n')
    first = True
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        if first:
            p = cell.paragraphs[0]
            first = False
        else:
            p = cell.add_paragraph()
        if stripped.startswith('- ') or stripped.startswith('* '):
            r = p.add_run(f"\u2022  {stripped[2:]}")
            r.font.name = "Arial"
            r.font.size = Pt(10)
            r.font.color.rgb = _BRS_BODY
        else:
            _brs_add_inline(p, stripped)
            for run in p.runs:
                run.font.name = "Arial"
                run.font.size = Pt(10)
                run.font.color.rgb = _BRS_BODY

    doc.add_paragraph()


def _brs_callout_box(doc, title_text, body_text, title_color=None):
    """
    Single-cell table callout box with colored title.
    title_color: 'critical' (#E07020), 'prohibition' (#C00000), or default (#1B3A5C).
    """
    if title_color == 'critical':
        tc_rgb = _BRS_ORANGE
    elif title_color == 'prohibition':
        tc_rgb = _BRS_RED
    else:
        tc_rgb = _BRS_BLUE

    table = doc.add_table(rows=1, cols=1)
    table.autofit = True
    cell = table.cell(0, 0)
    _set_cell_shading(cell, _BRS_CALLOUT_BG)

    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcMar = OxmlElement('w:tcMar')
    for edge in ('top', 'bottom', 'start', 'end'):
        mar = OxmlElement(f'w:{edge}')
        mar.set(qn('w:w'), '160')
        mar.set(qn('w:type'), 'dxa')
        tcMar.append(mar)
    tcPr.append(tcMar)

    # Title
    p = cell.paragraphs[0]
    r = p.add_run(title_text)
    r.font.name = "Arial"
    r.font.size = Pt(11)
    r.font.color.rgb = tc_rgb
    r.bold = True

    # Body lines
    for line in body_text.strip().split('\n'):
        stripped = line.strip()
        if not stripped:
            continue
        bp = cell.add_paragraph()
        _brs_add_inline(bp, stripped)
        for run in bp.runs:
            run.font.name = "Arial"
            run.font.size = Pt(10)
            run.font.color.rgb = _BRS_BODY

    doc.add_paragraph()


def _brs_render_body(doc, content, font_name="Arial"):
    """Render markdown content into Brisen house-style Word elements."""
    lines = content.split('\n')
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if not stripped:
            i += 1
            continue

        # Horizontal rules
        if stripped in ('---', '***', '___'):
            rule = doc.add_paragraph()
            _add_thin_border(rule, "CCCCCC")
            i += 1
            continue

        # Markdown table detection
        if '|' in stripped and i + 1 < len(lines):
            table_lines = []
            j = i
            while j < len(lines) and '|' in lines[j].strip():
                table_lines.append(lines[j].strip())
                j += 1
            if len(table_lines) >= 2:
                _brs_render_table(doc, table_lines, font_name)
                i = j
                continue

        # H1: Arial 14pt bold #1B3A5C
        if stripped.startswith('# ') and not stripped.startswith('## '):
            h = doc.add_heading(stripped[2:], level=1)
            for run in h.runs:
                run.font.name = font_name
                run.font.size = Pt(14)
                run.font.color.rgb = _BRS_BLUE
                run.bold = True
            i += 1
            continue

        # H2: Arial 12pt bold #1B3A5C
        if stripped.startswith('## ') and not stripped.startswith('### '):
            h = doc.add_heading(stripped[3:], level=2)
            for run in h.runs:
                run.font.name = font_name
                run.font.size = Pt(12)
                run.font.color.rgb = _BRS_BLUE
                run.bold = True
            i += 1
            continue

        # H3: Arial 11pt bold #1B3A5C
        if stripped.startswith('### '):
            h = doc.add_heading(stripped[4:], level=3)
            for run in h.runs:
                run.font.name = font_name
                run.font.size = Pt(11)
                run.font.color.rgb = _BRS_BLUE
                run.bold = True
            i += 1
            continue

        # H4
        if stripped.startswith('#### '):
            h = doc.add_heading(stripped[5:], level=4)
            for run in h.runs:
                run.font.name = font_name
                run.font.size = Pt(10)
                run.font.color.rgb = _BRS_BLUE
            i += 1
            continue

        # Blockquote → callout box (left border)
        if stripped.startswith('> '):
            quote_lines = []
            while i < len(lines) and lines[i].strip().startswith('> '):
                quote_lines.append(lines[i].strip()[2:])
                i += 1
            _brs_callout_box(doc, "", '\n'.join(quote_lines))
            continue

        # Bullet
        if stripped.startswith('- ') or stripped.startswith('* '):
            p = doc.add_paragraph(style='List Bullet')
            _brs_add_inline(p, stripped[2:])
            for run in p.runs:
                run.font.name = font_name
                run.font.size = Pt(10)
            i += 1
            continue

        # Sub-bullet
        if line.startswith('  - ') or line.startswith('  * ') or \
           line.startswith('    - ') or line.startswith('    * '):
            p = doc.add_paragraph(style='List Bullet 2')
            text = stripped.lstrip('-* ')
            _brs_add_inline(p, text)
            for run in p.runs:
                run.font.name = font_name
                run.font.size = Pt(10)
            i += 1
            continue

        # Numbered list
        if re.match(r'^\d+\.\s', stripped):
            text = re.sub(r'^\d+\.\s', '', stripped)
            p = doc.add_paragraph(style='List Number')
            _brs_add_inline(p, text)
            for run in p.runs:
                run.font.name = font_name
                run.font.size = Pt(10)
            i += 1
            continue

        # Q&A detection: line starting with Q followed by number
        if re.match(r'^Q\d+:', stripped):
            p = doc.add_paragraph()
            _brs_add_inline(p, stripped)
            for run in p.runs:
                run.font.name = font_name
                run.font.size = Pt(11)
                run.bold = True
            i += 1
            continue

        # Regular paragraph
        para = doc.add_paragraph()
        _brs_add_inline(para, stripped)
        for run in para.runs:
            run.font.name = font_name
            run.font.size = Pt(10)
        i += 1


def _brs_add_inline(paragraph, text):
    """Parse **bold**, *italic*, and `code` inline markdown."""
    parts = re.split(r'(\*\*.*?\*\*|\*.*?\*|`.*?`)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            run = paragraph.add_run(part[2:-2])
            run.bold = True
        elif part.startswith('*') and part.endswith('*'):
            run = paragraph.add_run(part[1:-1])
            run.italic = True
        elif part.startswith('`') and part.endswith('`'):
            run = paragraph.add_run(part[1:-1])
            run.font.name = "Consolas"
            run.font.size = Pt(9)
            run.font.color.rgb = _BRS_BLUE
        else:
            paragraph.add_run(part)


def _brs_render_table(doc, table_lines, font_name="Arial"):
    """Render markdown table matching 007 style — dark blue header, white text."""
    headers = []
    rows = []

    for line in table_lines:
        cells = [c.strip() for c in line.split('|') if c.strip()]
        if all(re.match(r'^[-:]+$', c) for c in cells):
            continue
        if not headers:
            headers = cells
        else:
            rows.append(cells)

    if not headers:
        return

    ncols = len(headers)
    table = doc.add_table(rows=1 + len(rows), cols=ncols)
    table.autofit = True

    # Header row: #1B3A5C bg, white bold text
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.paragraphs[0]
        r = p.add_run(h)
        r.font.name = font_name
        r.font.size = Pt(10)
        r.font.color.rgb = _BRS_WHITE
        r.bold = True
        _set_cell_shading(cell, _BRS_TABLE_HDR)

    # Data rows with alternating shading
    for i, row_data in enumerate(rows):
        for j, val in enumerate(row_data):
            if j >= ncols:
                break
            cell = table.cell(i + 1, j)
            cell.text = ""
            p = cell.paragraphs[0]
            r = p.add_run(val)
            r.font.name = font_name
            r.font.size = Pt(10)
            r.font.color.rgb = _BRS_BODY
            if i % 2 == 1:
                _set_cell_shading(cell, _BRS_TABLE_ALT)

    doc.add_paragraph()
